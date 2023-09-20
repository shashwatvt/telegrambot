import csv
import email
import os
import asyncio
import aiogram.types as types
import docx
import openpyxl
import pptx
import pypdf
from bs4 import BeautifulSoup
class FileTranscript:
    def __init__(self):
        self.VALID_EXTENSIONS = [
            "txt",
            "rtf",
            "md",
            "html",
            "xml",
            "csv",
            "json",
            "js",
            "css",
            "py",
            "java",
            "c",
            "cpp",
            "php",
            "rb",
            "swift",
            "sql",
            "sh",
            "bat",
            "ps1",
            "ini",
            "cfg",
            "conf",
            "log",
            "svg",
            "epub",
            "mobi",
            "tex",
            "docx",
            "odt",
            "xlsx",
            "ods",
            "pptx",
            "odp",
            "eml",
            "htaccess",
            "nginx.conf",
            "pdf",
        ]
    async def read_document(self, filename):
        try:
            extension = filename.split(".")[-1]
            contents = ''
            if extension not in self.VALID_EXTENSIONS:
                contents = "Invalid document file"
            if extension == "pdf":
                with open(filename, "rb") as f:
                    pdf_reader = pypdf.PdfReader(f)
                    num_pages = len(pdf_reader.pages)
                    for page_num in range(num_pages):
                        page_obj = pdf_reader.pages[page_num]
                        page_text = page_obj.extract_text()
                        contents += page_text
            elif extension == "docx":
                doc = docx.Document(filename)
                contents = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif extension in ["xlsx", "ods"]:
                workbook = openpyxl.load_workbook(filename, read_only=True)
                sheet = workbook.active
                for row in sheet.iter_rows(values_only=True):
                    contents += "\t".join([str(cell_value) for cell_value in row]) + "\n"
            elif extension in ["pptx", "odp"]:
                presentation = pptx.Presentation(filename)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            contents += shape.text + "\n"
            elif extension == "eml":
                with open(filename, "r") as f:
                    msg = email.message_from_file(f)
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            contents += part.get_payload()
            elif extension in ["html", "xml"]:
                with open(filename, "r") as f:
                    soup = BeautifulSoup(f, "html.parser")
                    contents = soup.get_text()
            elif extension == "csv":
                with open(filename, "r") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        contents += "\t".join(row) + "\n"
            else:
                with open(filename, "r") as f:
                    contents = f.read()
        except Exception as e:
            contents = f"Error during file download: {str(e)}"
        return contents
    async def download_file(self, bot, message: types.Message):
        try:
            if message.document is not None:
                file = message.document
                file_extension = (
                    file.file_name.split(".")[-1] if file.file_name is not None else ""
                )
            else:
                return None
            file_path = f"{file.file_id}.{file_extension}"
            file_dir = "downloaded_files"
            os.makedirs(file_dir, exist_ok=True)
            full_file_path = os.path.join(file_dir, file_path)
            await bot.download(file=file , destination=full_file_path)
            return full_file_path
        except Exception as e:
            print(f"Error during file download: {str(e)}")
            return None